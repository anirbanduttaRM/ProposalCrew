# ---- File Reading Logic ----
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import litellm
from model_config import get_model  # Import LLM configuration
import requests

# ---- RFP File Reading Functions ----

def read_rfp_txt(file_path):
    """Read RFP content from a .txt file."""
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    return content

def read_rfp_pdf(file_path):
    """Read RFP content from a PDF file."""
    reader = PdfReader(file_path)
    content = ""
    for page in reader.pages:
        content += page.extract_text() or ''  # Handle potential None values
    return content

def read_rfp_docx(file_path):
    """Read RFP content from a DOCX file."""
    doc = Document(file_path)
    content = "\n".join([para.text for para in doc.paragraphs])
    return content


# ---- Export to PowerPoint ----

def export_to_ppt(proposal, output_file="Proposal.pptx"):
    """
    Export the generated proposal to a PowerPoint presentation.
    """
    
    prs = Presentation()

    # ✅ Add Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Software Development Proposal"
    slide.placeholders[1].text = "Generated by CrewAI"

    # ✅ Add content slides
    for section in proposal.split("\n\n"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Section"
        
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        text_frame = text_box.text_frame

        paragraphs = section.split("\n")
        for para in paragraphs:
            p = text_frame.add_paragraph()
            p.text = para
            p.space_after = Inches(0.1)

    # ✅ Save the presentation
    prs.save(output_file)
    print(f"✅ Proposal exported as {output_file}")
